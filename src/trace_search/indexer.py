"""Document indexer for wiki knowledge base."""

import json
import logging
import re
import shutil
from collections.abc import Callable
from functools import lru_cache
from pathlib import Path

import bm25s
import chromadb
import Stemmer
from chromadb.errors import NotFoundError
from chromadb.config import Settings as ChromaSettings

from trace_search.config import settings
from trace_search.embeddings import EmbeddingBackend, build_embedding_backend

logger = logging.getLogger(__name__)

__all__ = [
    "SUPPORTED_EXTENSIONS",
    "SUPPORTED_EXTENSIONS_ORDERED",
    "WikiIndexer",
    "TokenCounter",
    "extract_title",
    "extract_content",
    "extract_pdf_content",
    "extract_docx_content",
    "extract_pptx_content",
    "extract_csv_content",
    "extract_code_content",
    "extract_notebook_content",
    "get_default_index_root",
    "should_exclude_path",
    "chunk_by_headings",
    "chunk_by_paragraphs",
    "create_contextual_chunk",
]

# Supported file extensions for indexing
SUPPORTED_EXTENSIONS_ORDERED = (
    # Documents
    ".md",
    ".pdf",
    ".docx",
    ".pptx",
    ".csv",
    # Code
    ".sql",
    ".py",
    ".yml",
    ".yaml",
    ".ipynb",
    # TypeScript
    ".ts",
    ".tsx",
)
SUPPORTED_EXTENSIONS = set(SUPPORTED_EXTENSIONS_ORDERED)


class TokenCounter:
    """Token counter using the embedding model's tokenizer.

    Use _get_token_counter() to get a cached singleton instance.
    """

    _tokenizer = None

    def _ensure_tokenizer(self) -> None:
        """Lazy-load tokenizer on first use."""
        if self._tokenizer is None:
            from transformers import AutoTokenizer

            self._tokenizer = AutoTokenizer.from_pretrained(
                settings.tokenizer_model,
                use_fast=True,
            )

    def count_tokens(self, text: str) -> int:
        """Count tokens in text using the embedding model's tokenizer."""
        self._ensure_tokenizer()
        return len(self._tokenizer.encode(text, add_special_tokens=False))

    def truncate_to_tokens(self, text: str, max_tokens: int) -> str:
        """Truncate text to fit within max_tokens."""
        self._ensure_tokenizer()
        tokens = self._tokenizer.encode(text, add_special_tokens=False)
        if len(tokens) <= max_tokens:
            return text
        truncated_tokens = tokens[:max_tokens]
        return self._tokenizer.decode(truncated_tokens, skip_special_tokens=True)

    def get_last_n_tokens(self, text: str, n: int) -> str:
        """Get the last n tokens of text as a string."""
        self._ensure_tokenizer()
        tokens = self._tokenizer.encode(text, add_special_tokens=False)
        if len(tokens) <= n:
            return text
        last_tokens = tokens[-n:]
        return self._tokenizer.decode(last_tokens, skip_special_tokens=True)

    def split_to_token_windows(
        self,
        text: str,
        *,
        max_tokens: int,
        overlap_tokens: int,
    ) -> list[str]:
        """Split text into token windows without dropping content."""
        self._ensure_tokenizer()
        tokens = self._tokenizer.encode(text, add_special_tokens=False)
        if len(tokens) <= max_tokens:
            return [text]

        step = (
            max_tokens - overlap_tokens
            if 0 < overlap_tokens < max_tokens
            else max_tokens
        )
        chunks = []
        for start in range(0, len(tokens), step):
            window = tokens[start : start + max_tokens]
            if not window:
                break
            chunks.append(self._tokenizer.decode(window, skip_special_tokens=True))
            if start + max_tokens >= len(tokens):
                break
        return chunks


@lru_cache(maxsize=1)
def _get_token_counter() -> TokenCounter:
    """Get singleton TokenCounter instance (lazy-loaded)."""
    return TokenCounter()


def _relative_parts(kb_path: Path, path: Path) -> tuple[str, ...]:
    """Return path parts relative to the KB root, resolving only when needed."""
    try:
        return path.relative_to(kb_path).parts
    except ValueError:
        return path.resolve().relative_to(kb_path.resolve()).parts


def _is_within_root(path: Path, root: Path) -> bool:
    """Return whether the resolved path stays under the resolved root."""
    try:
        return path.resolve().is_relative_to(root.resolve())
    except OSError:
        return False


def should_exclude_path(
    path: Path,
    kb_path: Path,
    exclude_patterns: list[str] | None = None,
) -> bool:
    """Return whether a KB-relative path should be skipped."""
    if not _is_within_root(path, kb_path):
        return True
    exclude = set(
        exclude_patterns if exclude_patterns is not None else settings.exclude_patterns_list
    )
    return any(
        part.startswith(".") or part in exclude
        for part in _relative_parts(kb_path, path)
    )


def get_default_index_root(
    kb_path: Path,
    index_root: Path | None = None,
    collection_name: str | None = None,
) -> Path:
    """Resolve the root directory that contains model-specific indexes."""
    if index_root is None:
        return kb_path / ".mcp-search" / "indexes"
    if collection_name:
        return index_root / collection_name
    return index_root


def extract_title(content: str, path: Path) -> str:
    """Extract title from file with fallback strategies."""
    ext = path.suffix.lower()

    # Code files: use filename with extension (e.g., "schema_parser.py")
    if ext in {".sql", ".py", ".ts", ".tsx", ".yml", ".yaml", ".ipynb"}:
        return path.name

    # Markdown: extract from first heading
    if ext == ".md":
        match = re.search(r"^#\s+(.+)$", content, re.MULTILINE)
        if match:
            return match.group(1).strip()
        return path.stem

    # Documents: find first substantial line
    if ext in {".pdf", ".docx", ".pptx", ".csv"}:
        for line in content.split("\n")[:10]:
            line = line.strip()
            if line and 10 < len(line) < 200 and not line.isupper():
                return line

    return path.stem


def extract_pdf_content(path: Path) -> str:
    """Extract text from PDF using PyMuPDF (better quality)."""
    import fitz  # PyMuPDF

    with fitz.open(path) as doc:
        pages = []
        for page in doc:
            text = page.get_text("text")
            if text.strip():
                pages.append(text)
    return "\n\n".join(pages)


def extract_docx_content(path: Path) -> str:
    """Extract text from Word document (paragraphs + tables)."""
    from docx import Document

    doc = Document(path)
    texts = []

    for para in doc.paragraphs:
        if para.text.strip():
            texts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                texts.append(row_text)

    return "\n\n".join(texts)


def extract_pptx_content(path: Path) -> str:
    """Extract text from PowerPoint (text frames + tables + notes)."""
    from pptx import Presentation

    prs = Presentation(path)
    texts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text)
            elif shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        slide_texts.append(row_text)

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(f"[Notes: {notes_text}]")

        texts.append("\n".join(slide_texts))

    return "\n\n".join(texts)


def extract_csv_content(path: Path) -> str:
    """Extract and flatten CSV data into readable text."""
    import csv

    with open(path, encoding="utf-8") as f:
        reader = csv.DictReader(f)
        rows = []
        for row_num, row in enumerate(reader, 1):
            row_parts = [f"Row {row_num}:"]
            for key, value in row.items():
                if value:
                    value_str = str(value).strip()
                    if value_str:
                        row_parts.append(f"{key}: {value_str}")
            if len(row_parts) > 1:  # Has content beyond row number
                rows.append(" | ".join(row_parts))
        return "\n".join(rows)


def extract_code_content(path: Path) -> str:
    """Extract code file content (SQL, Python, TypeScript, YAML)."""
    return path.read_text(encoding="utf-8", errors="replace")


def extract_notebook_content(path: Path) -> str:
    """Extract Jupyter notebook content (code + markdown cells)."""
    try:
        nb = json.loads(path.read_text(encoding="utf-8"))
        cells = nb.get("cells", [])
        content_parts = []
        for i, cell in enumerate(cells):
            cell_type = cell.get("cell_type", "code")
            source = "".join(cell.get("source", []))
            if source.strip():
                content_parts.append(f"[{cell_type.upper()} CELL {i + 1}]\n{source}")
        return "\n\n".join(content_parts)
    except Exception as e:
        logger.warning("Failed to parse notebook %s: %s", path, e)
        return ""


def extract_content(path: Path) -> str:
    """Extract text content based on file type.

    Args:
        path: Path to the file to extract content from.

    Returns:
        Extracted text content.

    Raises:
        ValueError: If file type is not supported.
    """
    ext = path.suffix.lower()

    extractors: dict[str, Callable[[Path], str]] = {
        ".md": lambda p: p.read_text(encoding="utf-8"),
        ".pdf": extract_pdf_content,
        ".docx": extract_docx_content,
        ".pptx": extract_pptx_content,
        ".csv": extract_csv_content,
        ".sql": extract_code_content,
        ".py": extract_code_content,
        ".ts": extract_code_content,
        ".tsx": extract_code_content,
        ".yml": extract_code_content,
        ".yaml": extract_code_content,
        ".ipynb": extract_notebook_content,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return extractor(path)


def _get_size(text: str, use_tokens: bool) -> int:
    """Get size of text in characters or tokens."""
    if use_tokens:
        return _get_token_counter().count_tokens(text)
    return len(text)


def _get_overlap_text(chunk: str, overlap_size: int, use_tokens: bool) -> str:
    """Extract overlap text from end of chunk."""
    if overlap_size <= 0:
        return ""
    if use_tokens:
        return _get_token_counter().get_last_n_tokens(chunk, overlap_size)
    return chunk[-overlap_size:] if len(chunk) > overlap_size else chunk


def _split_oversized_text(
    text: str,
    *,
    use_tokens: bool,
    max_size: int,
    overlap_size: int,
) -> list[str]:
    """Split one oversized text block into windows without losing content."""
    if _get_size(text, use_tokens) <= max_size:
        return [text]

    if use_tokens:
        return _get_token_counter().split_to_token_windows(
            text,
            max_tokens=max_size,
            overlap_tokens=overlap_size,
        )

    step = max_size - overlap_size if 0 < overlap_size < max_size else max_size
    return [text[start : start + max_size] for start in range(0, len(text), step)]


def _prepend_overlap_if_fits(
    overlap: str,
    chunk: str,
    *,
    use_tokens: bool,
    max_size: int,
) -> str:
    """Prefix overlap only when it does not force truncating the new chunk."""
    if not overlap:
        return chunk
    combined = f"{overlap}\n\n{chunk}"
    if _get_size(combined, use_tokens) <= max_size:
        return combined
    return chunk


def _merge_tiny_chunks(
    chunks: list[str],
    *,
    use_tokens: bool,
    max_size: int,
) -> list[str]:
    """Merge tiny chunks into neighbors to preserve context and reduce fragmentation."""
    if len(chunks) <= 1:
        return chunks

    if use_tokens:
        tiny_threshold = max(5, int(max_size * 0.1))
    else:
        tiny_threshold = 200 if max_size >= 200 else max(10, int(max_size * 0.5))

    merged: list[str] = []
    i = 0
    while i < len(chunks):
        chunk = chunks[i]
        if _get_size(chunk, use_tokens) >= tiny_threshold:
            merged.append(chunk)
            i += 1
            continue

        # Prefer merging with next chunk to keep heading context with its section.
        if i + 1 < len(chunks):
            combined_next = f"{chunk}\n\n{chunks[i + 1]}"
            if _get_size(combined_next, use_tokens) <= max_size:
                merged.append(combined_next)
                i += 2
                continue

        # Fall back to previous chunk if it fits.
        if merged:
            combined_prev = f"{merged[-1]}\n\n{chunk}"
            if _get_size(combined_prev, use_tokens) <= max_size:
                merged[-1] = combined_prev
                i += 1
                continue

        merged.append(chunk)
        i += 1

    return merged


def _extract_heading_context(section: str) -> str:
    """Extract contiguous markdown headings from the start of a section."""
    lines = []
    for line in section.splitlines():
        if line.strip().startswith("#"):
            lines.append(line.strip())
            continue
        if line.strip() == "":
            continue
        break
    return "\n".join(lines).strip()


def _prefix_heading_context(
    chunks: list[str],
    context: str,
    *,
    use_tokens: bool,
    max_size: int,
) -> list[str]:
    """Prefix heading context to chunks where it fits, to preserve section meaning."""
    if not context or len(chunks) <= 1:
        return chunks

    prefixed: list[str] = [chunks[0]]
    for chunk in chunks[1:]:
        if chunk.lstrip().startswith(context):
            prefixed.append(chunk)
            continue
        combined = f"{context}\n\n{chunk}"
        if _get_size(combined, use_tokens) <= max_size:
            prefixed.append(combined)
        else:
            prefixed.append(chunk)
    return prefixed


def _get_heading_level(section: str) -> int | None:
    """Return heading level from the first heading line in a section, if any."""
    for line in section.splitlines():
        stripped = line.lstrip()
        if not stripped:
            continue
        if stripped.startswith("#"):
            match = re.match(r"^(#{1,6})\s", stripped)
            if match:
                return len(match.group(1))
        break
    return None


def _heading_level_changed(current_level: int | None, next_level: int | None) -> bool:
    """Decide whether to split chunks when heading levels change."""
    if current_level is None and next_level is None:
        return False
    if current_level is None or next_level is None:
        return True
    return current_level != next_level


def _resolve_chunking_params(
    use_tokens: bool | None,
    max_chunk_chars: int | None,
    max_chunk_tokens: int | None,
    overlap_chars: int | None,
    overlap_tokens: int | None,
    enable_overlap: bool | None,
) -> tuple[bool, int, int]:
    """Resolve chunking parameters with defaults from settings.

    Returns:
        Tuple of (use_tokens, max_size, overlap_size)
    """
    if use_tokens is None:
        use_tokens = settings.use_token_based_chunking
    if max_chunk_chars is None:
        max_chunk_chars = settings.char_chunk_size
    if max_chunk_tokens is None:
        max_chunk_tokens = settings.token_chunk_size
    if overlap_chars is None:
        overlap_chars = settings.char_overlap_size
    if overlap_tokens is None:
        overlap_tokens = settings.token_overlap_size
    if enable_overlap is None:
        enable_overlap = settings.enable_chunk_overlap

    max_size = max_chunk_tokens if use_tokens else max_chunk_chars
    overlap_size = overlap_tokens if use_tokens else overlap_chars
    if not enable_overlap:
        overlap_size = 0

    return use_tokens, max_size, overlap_size


def chunk_by_headings(
    content: str,
    max_chunk_chars: int | None = None,
    *,
    use_tokens: bool | None = None,
    max_chunk_tokens: int | None = None,
    overlap_chars: int | None = None,
    overlap_tokens: int | None = None,
    enable_overlap: bool | None = None,
) -> list[str]:
    """Split content by markdown headings, respecting size limits.

    Args:
        content: Markdown content to chunk.
        max_chunk_chars: Max characters per chunk.
        use_tokens: Use token-based sizing.
        max_chunk_tokens: Max tokens per chunk (token mode).
        overlap_chars: Character overlap between chunks.
        overlap_tokens: Token overlap between chunks.
        enable_overlap: Enable overlap.

    Returns:
        List of content chunks.
    """
    use_tokens, max_size, overlap_size = _resolve_chunking_params(
        use_tokens,
        max_chunk_chars,
        max_chunk_tokens,
        overlap_chars,
        overlap_tokens,
        enable_overlap,
    )

    # Split on headings (keep the heading with content)
    sections = re.split(r"(?=^#{1,3}\s)", content, flags=re.MULTILINE)
    sections = [s.strip() for s in sections if s.strip()]

    chunks: list[str] = []
    current_chunk = ""
    pending_overlap = ""  # overlap from previous chunk
    current_level: int | None = None

    for section in sections:
        section_level = _get_heading_level(section)
        section_size = _get_size(section, use_tokens)
        current_size = _get_size(current_chunk, use_tokens)

        if current_chunk and _heading_level_changed(current_level, section_level):
            chunks.append(current_chunk)
            current_chunk = ""
            current_level = None
            pending_overlap = ""
            current_size = 0

        if current_size + section_size <= max_size:
            current_chunk += "\n\n" + section if current_chunk else section
            if current_chunk == section:
                current_level = section_level
        else:
            if current_chunk:
                chunks.append(current_chunk)
                pending_overlap = _get_overlap_text(
                    current_chunk, overlap_size, use_tokens
                )
                current_level = None

            # If section itself is too large, split it further
            if section_size > max_size:
                sub_chunks = chunk_by_paragraphs(
                    section,
                    max_chunk_chars=max_chunk_chars,
                    use_tokens=use_tokens,
                    max_chunk_tokens=max_chunk_tokens,
                    overlap_chars=overlap_chars,
                    overlap_tokens=overlap_tokens,
                    enable_overlap=enable_overlap,
                )
                heading_context = _extract_heading_context(section)
                sub_chunks = _prefix_heading_context(
                    sub_chunks,
                    heading_context,
                    use_tokens=use_tokens,
                    max_size=max_size,
                )
                if pending_overlap and sub_chunks:
                    sub_chunks[0] = _prepend_overlap_if_fits(
                        pending_overlap,
                        sub_chunks[0],
                        use_tokens=use_tokens,
                        max_size=max_size,
                    )
                chunks.extend(sub_chunks)
                if sub_chunks:
                    pending_overlap = _get_overlap_text(
                        sub_chunks[-1], overlap_size, use_tokens
                    )
                current_chunk = ""
                current_level = None
            else:
                current_chunk = _prepend_overlap_if_fits(
                    pending_overlap,
                    section,
                    use_tokens=use_tokens,
                    max_size=max_size,
                )
                pending_overlap = ""
                current_level = section_level

    if current_chunk:
        chunks.append(current_chunk)

    if not chunks:
        return [content]

    return _merge_tiny_chunks(
        chunks,
        use_tokens=use_tokens,
        max_size=max_size,
    )


def chunk_by_paragraphs(
    content: str,
    max_chunk_chars: int | None = None,
    *,
    use_tokens: bool | None = None,
    max_chunk_tokens: int | None = None,
    overlap_chars: int | None = None,
    overlap_tokens: int | None = None,
    enable_overlap: bool | None = None,
) -> list[str]:
    """Split content by paragraphs for large sections.

    Args:
        content: Content to chunk.
        max_chunk_chars: Max characters per chunk.
        use_tokens: Use token-based sizing.
        max_chunk_tokens: Max tokens per chunk (token mode).
        overlap_chars: Character overlap between chunks.
        overlap_tokens: Token overlap between chunks.
        enable_overlap: Enable overlap.

    Returns:
        List of content chunks.
    """
    use_tokens, max_size, overlap_size = _resolve_chunking_params(
        use_tokens,
        max_chunk_chars,
        max_chunk_tokens,
        overlap_chars,
        overlap_tokens,
        enable_overlap,
    )

    paragraphs = content.split("\n\n")
    chunks: list[str] = []
    current_chunk = ""
    pending_overlap = ""

    for para in paragraphs:
        if not para:
            continue
        para_size = _get_size(para, use_tokens)

        if para_size > max_size:
            if current_chunk:
                chunks.append(current_chunk)
                pending_overlap = _get_overlap_text(
                    current_chunk, overlap_size, use_tokens
                )
                current_chunk = ""

            sub_chunks = _split_oversized_text(
                para,
                use_tokens=use_tokens,
                max_size=max_size,
                overlap_size=overlap_size,
            )
            if pending_overlap and sub_chunks:
                sub_chunks[0] = _prepend_overlap_if_fits(
                    pending_overlap,
                    sub_chunks[0],
                    use_tokens=use_tokens,
                    max_size=max_size,
                )
                pending_overlap = ""
            chunks.extend(sub_chunks)
            if sub_chunks:
                pending_overlap = _get_overlap_text(
                    sub_chunks[-1], overlap_size, use_tokens
                )
            continue

        if current_chunk:
            candidate = f"{current_chunk}\n\n{para}"
        else:
            candidate = _prepend_overlap_if_fits(
                pending_overlap,
                para,
                use_tokens=use_tokens,
                max_size=max_size,
            )

        if _get_size(candidate, use_tokens) <= max_size:
            current_chunk = candidate
            if pending_overlap:
                pending_overlap = ""
        else:
            if current_chunk:
                chunks.append(current_chunk)
                pending_overlap = _get_overlap_text(
                    current_chunk, overlap_size, use_tokens
                )

            # Start new chunk with overlap if applicable
            current_chunk = _prepend_overlap_if_fits(
                pending_overlap,
                para,
                use_tokens=use_tokens,
                max_size=max_size,
            )
            pending_overlap = ""

    if current_chunk:
        chunks.append(current_chunk)

    return chunks


def create_contextual_chunk(title: str, folder: str, chunk: str) -> str:
    """Add document context to chunk (Anthropic's contextual retrieval pattern)."""
    return f"Document: {title}\nFolder: {folder}\n\n{chunk}"


def extract_breadcrumb(chunk: str, title: str) -> str:
    """Extract a readable heading breadcrumb from a chunk."""
    headings = []
    for line in chunk.splitlines():
        stripped = line.strip()
        match = re.match(r"^(#{1,6})\s+(.+)$", stripped)
        if match:
            headings.append(match.group(2).strip())
    if headings:
        return " > ".join(headings[-3:])
    return title


class WikiIndexer:
    """Index wiki documents into ChromaDB and BM25 for search."""

    def __init__(
        self,
        kb_path: str | Path | None = None,
        chroma_path: str | Path | None = None,
        bm25_path: str | Path | None = None,
        backend: EmbeddingBackend | None = None,
    ):
        """Initialize wiki indexer.

        Args:
            kb_path: Path to knowledge base. Uses KB_PATH env var if None.
            chroma_path: Path for ChromaDB storage. Auto-generated if None.
            bm25_path: Path for BM25 index storage. Auto-generated if None.
            backend: Pre-loaded embedding backend to share across indexers.
        """
        self.kb_path = Path(kb_path) if kb_path else settings.resolved_kb_path

        # Model-specific index paths to prevent dimension mismatch
        model_slug = settings.model_slug

        # Use explicit INDEX_PATH when set; otherwise keep indexes in the
        # documented per-KB .mcp-search/indexes directory.
        index_base = get_default_index_root(self.kb_path, settings.index_path)

        if chroma_path:
            self.chroma_path = Path(chroma_path)
        elif settings.chroma_path:
            self.chroma_path = Path(settings.chroma_path)
        else:
            self.chroma_path = index_base / f".chroma_db_{model_slug}"
        if bm25_path:
            self.bm25_path = Path(bm25_path)
        else:
            self.bm25_path = index_base / f".bm25_index_{model_slug}"

        logger.info(
            "Embedding model: %s (dims=%d)",
            settings.embedding_model,
            settings.embedding_dims,
        )
        logger.info("ChromaDB path: %s", self.chroma_path)
        logger.info("BM25 path: %s", self.bm25_path)

        self.client = chromadb.PersistentClient(
            path=str(self.chroma_path),
            settings=ChromaSettings(anonymized_telemetry=False),
        )

        self.backend: EmbeddingBackend = backend or build_embedding_backend()

        self.collection = self.client.get_or_create_collection(
            name="wiki_docs",
            metadata={"hnsw:space": "cosine"},
        )

        self._bm25: bm25s.BM25 | None = None
        self._bm25_corpus: list[dict] | None = None
        self._stemmer = Stemmer.Stemmer("english")

    def _get_relative_path(self, path: Path) -> str:
        return str(path.relative_to(self.kb_path))

    def _get_folder(self, path: Path) -> str:
        rel = path.relative_to(self.kb_path)
        parts = rel.parts
        return parts[0] if len(parts) > 1 else ""

    def _should_exclude(self, path: Path) -> bool:
        return should_exclude_path(path, self.kb_path)

    def _load_single_document(self, file_path: Path) -> dict | None:
        """Extract one supported file into a doc dict, or return None to skip."""
        ext = file_path.suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            return None
        if self._should_exclude(file_path):
            return None

        try:
            content = extract_content(file_path)
        except Exception as e:
            logger.warning("Failed to extract %s: %s", file_path, e)
            return None

        if not content.strip():
            return None

        stat = file_path.stat()
        return {
            "path": self._get_relative_path(file_path),
            "title": extract_title(content, file_path),
            "folder": self._get_folder(file_path),
            "extension": ext,
            "mtime": stat.st_mtime,
            "content": content,
        }

    def load_documents(self) -> list[dict]:
        """Load all supported files from knowledge base."""
        docs: list[dict] = []
        for file_path in self.kb_path.rglob("*"):
            doc = self._load_single_document(file_path)
            if doc is not None:
                docs.append(doc)
        docs.sort(key=lambda d: d["path"])
        return docs

    def _load_documents_subset(self, relative_paths: list[str]) -> list[dict]:
        """Load only the listed relative paths into doc dicts."""
        docs: list[dict] = []
        for rel in relative_paths:
            file_path = self.kb_path / rel
            if not file_path.is_file():
                continue
            doc = self._load_single_document(file_path)
            if doc is not None:
                docs.append(doc)
        docs.sort(key=lambda d: d["path"])
        return docs

    def _reconcile_index_state(self, force: bool) -> str:
        """Decide whether to do a full rebuild or an incremental update.

        Returns "force" when the indexes must be rebuilt from scratch (caller
        requested it, indexes are missing or empty, or metadata is missing or
        outdated). Returns "incremental" when the existing indexes are healthy
        enough to update in place; BM25 may need rebuilding from Chroma but
        chunks are preserved across the update.
        """
        from trace_search.index_metadata import (
            metadata_matches_active_model,
            read_index_metadata,
        )

        if force:
            return "force"

        chroma_count = self.collection.count()
        if chroma_count == 0:
            return "force"

        metadata = read_index_metadata(self.bm25_path.parent)
        if metadata is None:
            logger.info(
                "Index metadata missing or outdated; promoting to full rebuild"
            )
            return "force"

        if not metadata_matches_active_model(metadata):
            logger.info(
                "Index metadata does not match active embedding settings; "
                "promoting to full rebuild"
            )
            return "force"

        if self.bm25_path.exists():
            try:
                self._load_bm25()
            except Exception as e:
                logger.warning(
                    "BM25 load failed; will rebuild from Chroma during incremental: %s",
                    e,
                )
                self._bm25 = None
                self._bm25_corpus = None

        return "incremental"

    def _build_chunks(
        self, docs: list[dict]
    ) -> tuple[list[str], list[str], list[dict]]:
        """Convert documents into (chunks, ids, metadatas) ready for indexing."""
        all_chunks: list[str] = []
        all_ids: list[str] = []
        all_metadatas: list[dict] = []

        for doc in docs:
            chunks = chunk_by_headings(doc["content"])
            chunk_count = len(chunks)
            extension = doc.get("extension") or Path(doc["path"]).suffix.lower()
            source_mtime = float(doc.get("mtime") or 0.0)
            for i, chunk in enumerate(chunks):
                all_chunks.append(
                    create_contextual_chunk(doc["title"], doc["folder"], chunk)
                )
                all_ids.append(f"{doc['path']}::{i}")
                all_metadatas.append(
                    {
                        "path": doc["path"],
                        "title": doc["title"],
                        "folder": doc["folder"],
                        "chunk_index": i,
                        "chunk_count": chunk_count,
                        "breadcrumb": extract_breadcrumb(chunk, doc["title"]),
                        "extension": extension,
                        "source_mtime": source_mtime,
                    }
                )

        return all_chunks, all_ids, all_metadatas

    def _persist_chroma(
        self,
        ids: list[str],
        chunks: list[str],
        embeddings,
        metadatas: list[dict],
    ) -> None:
        """Write chunks and embeddings to ChromaDB in batches."""
        batch_size = 500
        for i in range(0, len(chunks), batch_size):
            end = min(i + batch_size, len(chunks))
            self.collection.add(
                ids=ids[i:end],
                documents=chunks[i:end],
                embeddings=embeddings[i:end].tolist(),
                metadatas=metadatas[i:end],
            )
            logger.info("Indexed %s/%s chunks (ChromaDB)", end, len(chunks))

    def _persist_bm25(self, chunks: list[str], metadatas: list[dict]) -> None:
        """Build and save the BM25 index plus metadata to disk."""
        logger.info("Building BM25 index...")
        corpus_tokens = bm25s.tokenize(
            chunks,
            stopwords="en",
            stemmer=self._stemmer,
        )
        self._bm25 = bm25s.BM25(k1=settings.bm25_k1, b=settings.bm25_b)
        self._bm25.index(corpus_tokens)
        self._bm25_corpus = metadatas

        self.bm25_path.mkdir(parents=True, exist_ok=True)
        self._bm25.save(str(self.bm25_path), corpus=chunks)

        metadata_path = self.bm25_path / "metadata.json"
        with open(metadata_path, "w") as f:
            json.dump(metadatas, f)

    def build_index(self, force: bool = False) -> int:
        """Build or update the search index. Returns number of chunks indexed.

        Default behavior is incremental: detect added, changed, and removed
        source files and apply only the necessary work. Pass `force=True` to
        drop both indexes and rebuild every file from scratch.
        """
        from trace_search.index_metadata import (
            build_index_metadata,
            categorize_source_changes,
            invalidate_index_metadata,
            read_index_metadata,
            utc_now_iso,
            write_index_metadata,
        )

        mode = self._reconcile_index_state(force)
        build_started_at = utc_now_iso()

        if mode == "force":
            return self._full_rebuild(build_started_at)

        metadata = read_index_metadata(self.bm25_path.parent)
        changes = categorize_source_changes(self.kb_path, metadata)

        bm25_healthy = (
            self.bm25_path.exists()
            and self._bm25 is not None
            and self._bm25_corpus is not None
        )
        if not changes.has_changes and bm25_healthy:
            logger.info(
                "Index up to date: %d unchanged files, %d chunks",
                len(changes.unchanged),
                self.collection.count(),
            )
            return self.collection.count()

        try:
            self._apply_incremental_changes(changes)
            self._rebuild_bm25_from_chroma()
        except Exception:
            logger.exception(
                "Incremental reindex failed; invalidating metadata so next "
                "reindex runs as a full rebuild"
            )
            invalidate_index_metadata(self.bm25_path.parent)
            raise

        chunk_count = self.collection.count()
        new_metadata = build_index_metadata(
            kb_path=self.kb_path,
            build_started_at=build_started_at,
            build_completed_at=utc_now_iso(),
            document_count=len(changes.inventory),
            chunk_count=chunk_count,
            source_files=changes.inventory,
        )
        write_index_metadata(self.bm25_path.parent, new_metadata)
        logger.info(
            "Incremental reindex: +%d added, ~%d changed, -%d removed, =%d unchanged; %d chunks total",
            len(changes.added),
            len(changes.changed),
            len(changes.removed),
            len(changes.unchanged),
            chunk_count,
        )
        return chunk_count

    def _full_rebuild(self, build_started_at: str) -> int:
        """Drop both indexes and reindex every visible file from scratch."""
        from trace_search.index_metadata import (
            build_index_metadata,
            collect_source_files,
            utc_now_iso,
            write_index_metadata,
        )

        self._clear_chroma_collection()
        self._clear_bm25_index()

        logger.info("Loading documents...")
        docs = self.load_documents()
        logger.info("Found %s documents", len(docs))

        all_chunks, all_ids, all_metadatas = self._build_chunks(docs)

        if not all_chunks:
            logger.info("No chunks to index")
            inventory = collect_source_files(self.kb_path)
            metadata = build_index_metadata(
                kb_path=self.kb_path,
                build_started_at=build_started_at,
                build_completed_at=utc_now_iso(),
                document_count=len(docs),
                chunk_count=0,
                source_files=inventory,
            )
            write_index_metadata(self.bm25_path.parent, metadata)
            return 0

        logger.info("Generating embeddings for %s chunks...", len(all_chunks))
        embeddings = self.backend.encode(all_chunks)

        self._persist_chroma(all_ids, all_chunks, embeddings, all_metadatas)
        self._persist_bm25(all_chunks, all_metadatas)

        chunk_count = self.collection.count()
        logger.info("Index complete: %s chunks (ChromaDB + BM25)", chunk_count)

        inventory = collect_source_files(self.kb_path)
        metadata = build_index_metadata(
            kb_path=self.kb_path,
            build_started_at=build_started_at,
            build_completed_at=utc_now_iso(),
            document_count=len(docs),
            chunk_count=chunk_count,
            source_files=inventory,
        )
        write_index_metadata(self.bm25_path.parent, metadata)
        return chunk_count

    def _apply_incremental_changes(self, changes) -> None:
        """Apply categorized file changes to the Chroma collection in place."""
        for path in changes.changed + changes.removed:
            self.collection.delete(where={"path": path})

        paths_to_index = changes.added + changes.changed
        if not paths_to_index:
            return

        docs = self._load_documents_subset(paths_to_index)
        chunks, ids, metadatas = self._build_chunks(docs)
        if not chunks:
            return

        embeddings = self.backend.encode(chunks)
        self._persist_chroma(ids, chunks, embeddings, metadatas)

    def _rebuild_bm25_from_chroma(self) -> None:
        """Rebuild BM25 from the current Chroma chunk inventory."""
        result = self.collection.get(include=["documents", "metadatas"])
        documents = list(result.get("documents") or [])
        metadatas = list(result.get("metadatas") or [])

        self._clear_bm25_index()

        if not documents:
            return

        self._persist_bm25(documents, metadatas)

    def _clear_chroma_collection(self) -> None:
        """Drop and recreate the Chroma collection atomically without materializing ids."""
        try:
            self.client.delete_collection("wiki_docs")
        except NotFoundError:
            logger.debug("Chroma collection did not exist before rebuild")
        except Exception as exc:
            logger.warning("Failed to delete Chroma collection before rebuild: %s", exc)
            raise
        self.collection = self.client.get_or_create_collection(
            name="wiki_docs",
            metadata={"hnsw:space": "cosine"},
        )

    def _clear_bm25_index(self) -> None:
        """Delete BM25 files and clear in-memory BM25 state."""
        self._bm25 = None
        self._bm25_corpus = None
        if self.bm25_path.exists():
            shutil.rmtree(self.bm25_path)

    def _load_bm25(self) -> None:
        """Load BM25 index from disk."""
        if self._bm25 is not None:
            return

        if not self.bm25_path.exists():
            return

        self._bm25 = bm25s.BM25.load(str(self.bm25_path), load_corpus=True)
        metadata_path = self.bm25_path / "metadata.json"
        if metadata_path.exists():
            with open(metadata_path) as f:
                self._bm25_corpus = json.load(f)

    @property
    def bm25(self) -> bm25s.BM25 | None:
        """Get BM25 index, loading from disk if needed."""
        if self._bm25 is None:
            self._load_bm25()
        return self._bm25

    @property
    def bm25_corpus(self) -> list[dict] | None:
        """Get BM25 corpus metadata."""
        if self._bm25_corpus is None:
            self._load_bm25()
        return self._bm25_corpus

    def get_stats(self) -> dict[str, object]:
        """Get index statistics."""
        bm25 = self.bm25
        bm25_corpus = self.bm25_corpus
        bm25_docs = len(bm25_corpus) if bm25_corpus else 0
        return {
            "total_chunks": self.collection.count(),
            "bm25_docs": bm25_docs,
            "kb_path": str(self.kb_path),
            "chroma_path": str(self.chroma_path),
            "bm25_path": str(self.bm25_path),
            "embedding_model": settings.embedding_model,
            "embedding_dims": settings.embedding_dims,
            "embedding_pooling": settings.embedding_pooling,
            "bm25_available": bm25 is not None,
            # Chunking configuration
            "chunking": {
                "use_tokens": settings.use_token_based_chunking,
                "enable_overlap": settings.enable_chunk_overlap,
                "char_chunk_size": settings.char_chunk_size,
                "char_overlap_size": settings.char_overlap_size,
                "token_chunk_size": settings.token_chunk_size,
                "token_overlap_size": settings.token_overlap_size,
            },
        }
