"""Document content extractors."""

from __future__ import annotations

import json
import logging
import re
from collections.abc import Callable
from pathlib import Path

logger = logging.getLogger(__name__)

# Supported file extensions for indexing
SUPPORTED_EXTENSIONS_ORDERED = (
    # Documents
    ".md",
    ".pdf",
    ".docx",
    ".pptx",
    ".csv",
    # Code
    ".sql",
    ".py",
    ".yml",
    ".yaml",
    ".ipynb",
    # TypeScript
    ".ts",
    ".tsx",
)
SUPPORTED_EXTENSIONS = set(SUPPORTED_EXTENSIONS_ORDERED)


def extract_title(content: str, path: Path) -> str:
    """Extract title from file with fallback strategies."""
    ext = path.suffix.lower()

    # Code files: use filename with extension (e.g., "schema_parser.py")
    if ext in {".sql", ".py", ".ts", ".tsx", ".yml", ".yaml", ".ipynb"}:
        return path.name

    # Markdown: extract from first heading
    if ext == ".md":
        match = re.search(r"^#\s+(.+)$", content, re.MULTILINE)
        if match:
            return match.group(1).strip()
        return path.stem

    # Documents: find first substantial line
    if ext in {".pdf", ".docx", ".pptx", ".csv"}:
        for line in content.split("\n")[:10]:
            line = line.strip()
            if line and 10 < len(line) < 200 and not line.isupper():
                return line

    return path.stem


def extract_pdf_content(path: Path) -> str:
    """Extract text from PDF using PyMuPDF (better quality)."""
    import fitz  # PyMuPDF

    with fitz.open(path) as doc:
        pages = []
        for page in doc:
            text = page.get_text("text")
            if text.strip():
                pages.append(text)
    return "\n\n".join(pages)


def extract_docx_content(path: Path) -> str:
    """Extract text from Word document (paragraphs + tables)."""
    from docx import Document

    doc = Document(path)
    texts = []

    for para in doc.paragraphs:
        if para.text.strip():
            texts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                texts.append(row_text)

    return "\n\n".join(texts)


def extract_pptx_content(path: Path) -> str:
    """Extract text from PowerPoint (text frames + tables + notes)."""
    from pptx import Presentation

    prs = Presentation(path)
    texts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text)
            elif shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        slide_texts.append(row_text)

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(f"[Notes: {notes_text}]")

        texts.append("\n".join(slide_texts))

    return "\n\n".join(texts)


def extract_csv_content(path: Path) -> str:
    """Extract and flatten CSV data into readable text."""
    import csv

    with open(path, encoding="utf-8") as f:
        reader = csv.DictReader(f)
        rows = []
        for row_num, row in enumerate(reader, 1):
            row_parts = [f"Row {row_num}:"]
            for key, value in row.items():
                if value:
                    value_str = str(value).strip()
                    if value_str:
                        row_parts.append(f"{key}: {value_str}")
            if len(row_parts) > 1:  # Has content beyond row number
                rows.append(" | ".join(row_parts))
        return "\n".join(rows)


def extract_code_content(path: Path) -> str:
    """Extract code file content (SQL, Python, TypeScript, YAML)."""
    return path.read_text(encoding="utf-8", errors="replace")


def extract_notebook_content(path: Path) -> str:
    """Extract Jupyter notebook content (code + markdown cells)."""
    try:
        nb = json.loads(path.read_text(encoding="utf-8"))
        cells = nb.get("cells", [])
        content_parts = []
        for i, cell in enumerate(cells):
            cell_type = cell.get("cell_type", "code")
            source = "".join(cell.get("source", []))
            if source.strip():
                content_parts.append(f"[{cell_type.upper()} CELL {i + 1}]\n{source}")
        return "\n\n".join(content_parts)
    except Exception as e:
        logger.warning("Failed to parse notebook %s: %s", path, e)
        return ""


def extract_content(path: Path) -> str:
    """Extract text content based on file type.

    Args:
        path: Path to the file to extract content from.

    Returns:
        Extracted text content.

    Raises:
        ValueError: If file type is not supported.
    """
    ext = path.suffix.lower()

    extractors: dict[str, Callable[[Path], str]] = {
        ".md": lambda p: p.read_text(encoding="utf-8"),
        ".pdf": extract_pdf_content,
        ".docx": extract_docx_content,
        ".pptx": extract_pptx_content,
        ".csv": extract_csv_content,
        ".sql": extract_code_content,
        ".py": extract_code_content,
        ".ts": extract_code_content,
        ".tsx": extract_code_content,
        ".yml": extract_code_content,
        ".yaml": extract_code_content,
        ".ipynb": extract_notebook_content,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return extractor(path)
